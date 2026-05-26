"""Build professional PowerPoint presentation for EWS project."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

FIGDIR = "presentation/figures"

# ── Colors ──
POLIMI_DARK = RGBColor(26, 58, 92)
POLIMI_BLUE = RGBColor(41, 128, 185)
POLIMI_RED = RGBColor(192, 57, 43)
POLIMI_GREEN = RGBColor(39, 174, 96)
POLIMI_GOLD = RGBColor(241, 196, 15)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(245, 245, 245)
DARK_GRAY = RGBColor(60, 60, 60)
MID_GRAY = RGBColor(140, 140, 140)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helper functions ──
def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_rich_text(slide, left, top, width, height, paragraphs):
    """paragraphs = list of (text, font_size, bold, color, spacing_after)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, fs, bold, color, sp_after) in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Calibri"
        if sp_after:
            p.space_after = Pt(sp_after)
    return txBox

def add_logo_corner(slide):
    """Add circular logo to top-right corner."""
    logo_path = "presentation/logo_circle.png"
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(12.2), Inches(0.2), Inches(0.8), Inches(0.8))

def add_section_bar(slide, text):
    """Add a thin colored bar at the top with section name."""
    bar = add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), POLIMI_DARK)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(5), Inches(0.4),
                 text, font_size=9, color=MID_GRAY, bold=False)

def add_slide_title(slide, title, subtitle=None):
    add_text_box(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.6),
                 title, font_size=28, bold=True, color=POLIMI_DARK)
    if subtitle:
        add_text_box(slide, Inches(0.7), Inches(0.95), Inches(11), Inches(0.4),
                     subtitle, font_size=14, color=MID_GRAY)

def add_footer(slide, page_num, total=15):
    add_text_box(slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
                 f"{page_num} / {total}", font_size=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)
    add_text_box(slide, Inches(0.7), Inches(7.1), Inches(5), Inches(0.3),
                 "Ercelli, Galli, Impenati, Mazzini — PoliMi GSoM", font_size=9, color=MID_GRAY)

def img(name):
    return os.path.join(FIGDIR, name)


# ════════════════════════════════════════════════
# SLIDE 1: COVER
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, WHITE)
# Dark top band
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(3.2), POLIMI_DARK)
# Accent line
add_shape(slide, Inches(0), Inches(3.2), prs.slide_width, Inches(0.06), POLIMI_BLUE)

# Full logo
if os.path.exists("presentation/logo_full.png"):
    slide.shapes.add_picture("presentation/logo_full.png", Inches(0.8), Inches(0.4), Inches(3.5))

# Title
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.8),
             "Early-Warning System for Risk-Off Regimes",
             font_size=36, bold=True, color=WHITE)
add_text_box(slide, Inches(0.8), Inches(2.1), Inches(11.5), Inches(0.6),
             "with Domain-Specific Allocation Routing",
             font_size=22, color=RGBColor(180, 210, 240))
add_text_box(slide, Inches(0.8), Inches(2.6), Inches(11.5), Inches(0.4),
             "Unsupervised Anomaly Detection for Regime-Aware Portfolio Allocation",
             font_size=14, color=RGBColor(160, 190, 220))

# Authors
add_text_box(slide, Inches(0.8), Inches(3.8), Inches(11.5), Inches(0.5),
             "Edoardo Ercelli  ·  Luca Galli  ·  Francesco Impenati  ·  Francesco Mazzini",
             font_size=18, bold=True, color=POLIMI_DARK)
add_text_box(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.4),
             "POLIMI Graduate School of Management — Politecnico di Milano",
             font_size=14, color=MID_GRAY)
add_text_box(slide, Inches(0.8), Inches(4.9), Inches(11.5), Inches(0.4),
             "New Technologies for Finance — May 2025",
             font_size=12, color=MID_GRAY)


# ════════════════════════════════════════════════
# SLIDE 2: PROBLEM STATEMENT
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_bar(slide, "PRESENTATION")
add_slide_title(slide, "Problem Statement & Motivation")
add_logo_corner(slide)
add_footer(slide, 2)

# Left column
paras = [
    ("The gap in traditional risk management", 16, True, POLIMI_DARK, 8),
    ("• Most early-warning systems treat risk-off as a binary decision: equity vs. cash", 12, False, DARK_GRAY, 4),
    ("• But crises are heterogeneous: in 2008 the USD was the safe haven (dollar funding squeeze); in 2020 it initially strengthened then weakened; in an inflationary stress regime, gold outperforms", 12, False, DARK_GRAY, 4),
    ("• No existing framework combines regime detection with allocation routing based on the nature of the stress", 12, False, DARK_GRAY, 10),
    ("Our contribution", 16, True, POLIMI_DARK, 6),
    ("An EWS that answers two questions simultaneously:", 12, False, DARK_GRAY, 4),
    ("1. WHEN to go risk-off (ensemble anomaly detection)", 12, True, POLIMI_BLUE, 3),
    ("2. WHERE to allocate (domain-specific macro routing)", 12, True, POLIMI_BLUE, 4),
]
add_rich_text(slide, Inches(0.7), Inches(1.3), Inches(5.8), Inches(5.5), paras)

# Right column — Investment thesis box
box = add_shape(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(3.2), LIGHT_GRAY)
paras_r = [
    ("Investment Thesis", 16, True, POLIMI_DARK, 8),
    ("Deploy 1.5× equity leverage during calm regimes to exploit the equity risk premium, and rely on the EWS to detect stress early enough to avoid drawdowns — maximizing the Sharpe ratio.", 12, False, DARK_GRAY, 8),
    ("The leverage is justified only if the detection system has sufficient recall: missing a crisis at 1.5× is costlier than at 1.0×.", 12, False, POLIMI_RED, 8),
    ("This creates a built-in accountability constraint on model performance.", 11, False, MID_GRAY, 0),
]
add_rich_text(slide, Inches(7.3), Inches(1.5), Inches(5), Inches(3), paras_r)

# References box
paras_ref = [
    ("Key references", 11, True, MID_GRAY, 4),
    ("Maggiori (2020) — dollar as global safe asset", 10, False, MID_GRAY, 2),
    ("Baur & Lucey (2010) — gold: hedge vs. safe haven", 10, False, MID_GRAY, 2),
    ("Avdjiev et al. (2019) — dollar-leverage-flows triangle", 10, False, MID_GRAY, 0),
]
add_rich_text(slide, Inches(7.3), Inches(4.7), Inches(5), Inches(2), paras_ref)


# ════════════════════════════════════════════════
# SLIDE 3: DATASET
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_bar(slide, "PRESENTATION")
add_slide_title(slide, "Dataset Overview")
add_logo_corner(slide)
add_footer(slide, 3)

paras_d = [
    ("Bloomberg weekly data — 2000 to 2021 — 1,111 observations", 14, True, POLIMI_DARK, 8),
    ("43 features: 7 MSCI equity, 3 FX, 4 commodities, 8 bond TR indices, 18 yields/rates, VIX, US Economic Surprise", 11, False, DARK_GRAY, 4),
    ("Target Y=1: risk-off week (21.3% prevalence, 237/1111 weeks)", 11, False, DARK_GRAY, 4),
    ("No missing values across entire dataset", 11, False, DARK_GRAY, 10),
    ("Proxy decisions (maintained throughout):", 12, True, POLIMI_DARK, 4),
    ("• MSCI World (absent) → MXUS as single-name proxy", 11, False, DARK_GRAY, 3),
    ("• Global Aggregate (absent) → LUACTRUU (US IG Corporate TR)", 11, False, DARK_GRAY, 3),
    ("• Credit spreads = log price-index ratios (not OAS yields)", 11, False, DARK_GRAY, 0),
]
add_rich_text(slide, Inches(0.7), Inches(1.3), Inches(5.5), Inches(5.5), paras_d)

# Figure
slide.shapes.add_picture(img("fig_y_distribution.png"), Inches(6.5), Inches(1.2), Inches(6.2))
add_text_box(slide, Inches(6.5), Inches(4.5), Inches(6.2), Inches(1.5),
             "Risk-off clusters in 2000–02, 2008–09, 2011–12, 2020. Mid-decade years have near-zero positives — this data property critically shapes cross-validation design.",
             font_size=10, color=MID_GRAY)


# ════════════════════════════════════════════════
# SLIDE 4: THREE MACRO DOMAINS
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_bar(slide, "BUSINESS LOGIC")
add_slide_title(slide, "Three Macroeconomic Domains for Routing")
add_logo_corner(slide)
add_footer(slide, 4)

# USD Domain box
add_shape(slide, Inches(0.5), Inches(1.3), Inches(3.9), Inches(5.5), RGBColor(235, 245, 255))
paras_usd = [
    ("USD Domain", 16, True, POLIMI_BLUE, 6),
    ("Dollar is safe haven when stress is global but external to the US, or during dollar funding squeeze (Avdjiev et al. 2019).", 10, False, DARK_GRAY, 8),
    ("Triggers (signed z-scores):", 11, True, POLIMI_DARK, 4),
    ("[+] LIBOR-3M spread Δ4w (funding stress)", 10, False, DARK_GRAY, 2),
    ("[+] DXY momentum Δ4w", 10, False, DARK_GRAY, 2),
    ("[+] VRP (implied > realized vol)", 10, False, DARK_GRAY, 2),
    ("[−] US 10Y Δ change (Treasury flight)", 10, False, DARK_GRAY, 2),
    ("[+] USA relative vs DM composite", 10, False, DARK_GRAY, 0),
]
add_rich_text(slide, Inches(0.7), Inches(1.5), Inches(3.5), Inches(5.2), paras_usd)

# Gold Domain box
add_shape(slide, Inches(4.7), Inches(1.3), Inches(3.9), Inches(5.5), RGBColor(255, 250, 230))
paras_oro = [
    ("Gold Domain", 16, True, RGBColor(200, 160, 0), 6),
    ("Gold when stress has monetary/inflationary component or the dollar itself is the problem (Baur & Lucey 2010; Erb & Harvey 2013).", 10, False, DARK_GRAY, 8),
    ("Triggers (signed z-scores):", 11, True, POLIMI_DARK, 4),
    ("[−] Real yield proxy Δ4w", 10, False, DARK_GRAY, 2),
    ("[−] DXY Δ4w (opposite sign!)", 10, False, DARK_GRAY, 2),
    ("[+] JPY strength (safe haven)", 10, False, DARK_GRAY, 2),
    ("[+] Equity-bond ρ 13w (diversif. failure)", 10, False, DARK_GRAY, 2),
    ("[+] Gold-oil ratio Δ4w", 10, False, DARK_GRAY, 0),
]
add_rich_text(slide, Inches(4.9), Inches(1.5), Inches(3.5), Inches(5.2), paras_oro)

# MBS Domain box
add_shape(slide, Inches(8.9), Inches(1.3), Inches(3.9), Inches(5.5), RGBColor(240, 240, 240))
paras_mbs = [
    ("MBS Domain", 16, True, MID_GRAY, 6),
    ("MBS is NOT a crisis refuge — it is defensive carry for moderate stress only (Diep et al. 2021).", 10, False, DARK_GRAY, 8),
    ("Rule-based activation:", 11, True, POLIMI_DARK, 4),
    ("• VIX ∈ [20, 28] (nervous, not panic)", 10, False, DARK_GRAY, 2),
    ("• Term spread > 0 (normal curve)", 10, False, DARK_GRAY, 2),
    ("• Drawdown ∈ [−12%, −5%]", 10, False, DARK_GRAY, 6),
    ("Blocked if:", 11, True, POLIMI_RED, 4),
    ("• VIX > 30 (acute stress)", 10, False, DARK_GRAY, 2),
    ("• Funding spread > p90", 10, False, DARK_GRAY, 2),
    ("• HY-IG widening > 50bps", 10, False, DARK_GRAY, 0),
]
add_rich_text(slide, Inches(9.1), Inches(1.5), Inches(3.5), Inches(5.2), paras_mbs)

# Bottom note
add_text_box(slide, Inches(0.5), Inches(6.9), Inches(12), Inches(0.3),
             "Note: DXY appears in both USD (+) and Gold (−) domains with opposite sign — same indicator, opposite economic reading. By design.",
             font_size=9, color=MID_GRAY)


# ════════════════════════════════════════════════
# SLIDE 5: DECISION MATRIX
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_bar(slide, "BUSINESS LOGIC")
add_slide_title(slide, "Routing Architecture & Decision Matrix")
add_logo_corner(slide)
add_footer(slide, 5)

paras_dm = [
    ("Allocation Decision Matrix", 14, True, POLIMI_DARK, 8),
    ("Risk-on (signal=0)    →  EQUITY (1.5× equity, −0.5× cash)", 12, False, POLIMI_GREEN, 3),
    ("Risk-off + USD > τ + DXY↑  →  CASH USD", 12, False, POLIMI_BLUE, 3),
    ("Risk-off + Gold > τ       →  GOLD", 12, False, RGBColor(200, 160, 0), 3),
    ("Risk-off + MBS active     →  MBS", 12, False, MID_GRAY, 3),
    ("Risk-off + default         →  CASH USD", 12, False, POLIMI_BLUE, 12),
    ("Design principles:", 14, True, POLIMI_DARK, 6),
    ("• Priority: USD → Gold → MBS (funding stress is most urgent)", 11, False, DARK_GRAY, 3),
    ("• Default = Cash: zero-cost, doesn't compound errors", 11, False, DARK_GRAY, 3),
    ("• TC differentiated: equity 5bps, cash 2bps, gold 8bps, MBS 20bps", 11, False, DARK_GRAY, 3),
    ("• Only 2 scalar thresholds optimized — rest is rule-based from macro rationale", 11, False, DARK_GRAY, 0),
]
add_rich_text(slide, Inches(0.7), Inches(1.3), Inches(5.8), Inches(5.5), paras_dm)

# Sub-score formula box
add_shape(slide, Inches(7), Inches(1.3), Inches(5.8), Inches(2.5), LIGHT_GRAY)
paras_ss = [
    ("Sub-Score Construction", 14, True, POLIMI_DARK, 6),
    ("S = (1/|D|) × Σ sign × (x − μ_dev) / σ_dev", 13, True, POLIMI_BLUE, 8),
    ("μ and σ estimated ONLY on development set (no leakage)", 11, False, POLIMI_RED, 6),
    ("Threshold optimization: grid search τ_usd × τ_gold,", 11, False, DARK_GRAY, 3),
    ("selecting by Calmar ratio weighted by fold duration", 11, False, DARK_GRAY, 0),
]
add_rich_text(slide, Inches(7.3), Inches(1.5), Inches(5.3), Inches(2.3), paras_ss)

# Calmar heatmap
slide.shapes.add_picture(img("fig_calmar_heatmap.png"), Inches(7), Inches(4.0), Inches(5.5))


# ════════════════════════════════════════════════
# SLIDE 6: FEATURE ENGINEERING
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_bar(slide, "PROCESS & EXECUTION")
add_slide_title(slide, "Feature Engineering: Stationarity & Cross-Asset Spreads")
add_logo_corner(slide)
add_footer(slide, 6)

paras_fe = [
    ("Stationarity transforms", 14, True, POLIMI_DARK, 6),
    ("Anomaly detection assumes stable distribution — non-stationary features violate this fundamentally.", 11, False, DARK_GRAY, 4),
    ("• Log-returns: prices (equity, FX, commodities, bond TR)", 11, False, DARK_GRAY, 2),
    ("• First differences: yields and rates", 11, False, DARK_GRAY, 2),
    ("• Level: VIX, ECSURPUS (already stationary)", 11, False, DARK_GRAY, 4),
    ("Result: 42/42 features pass ADF at 5%", 11, True, POLIMI_GREEN, 8),
    ("20 cross-asset spreads constructed:", 14, True, POLIMI_DARK, 4),
    ("• Term: US 10Y−2Y, DE 10Y−2Y, BTP−Bund, US−DE", 10, False, DARK_GRAY, 2),
    ("• Credit: HY-IG, EM (log price ratios, safe − risky)", 10, False, DARK_GRAY, 2),
    ("• VRP: VIX − realized vol (4w, annualized)", 10, False, DARK_GRAY, 2),
    ("• Equity-bond rotation, gold-oil ratio, JPY strength", 10, False, DARK_GRAY, 4),
    ("7 collinear features removed (|ρ| > 0.9) → 56 clean features", 11, True, POLIMI_DARK, 0),
]
add_rich_text(slide, Inches(0.7), Inches(1.3), Inches(5.5), Inches(5.8), paras_fe)

slide.shapes.add_picture(img("fig_routing_corr.png"), Inches(6.5), Inches(1.3), Inches(6.2))


# ════════════════════════════════════════════════
# SLIDE 7: WALK-FORWARD CV
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_bar(slide, "PROCESS & EXECUTION")
add_slide_title(slide, "Walk-Forward Cross-Validation with Purged Embargo")
add_logo_corner(slide)
add_footer(slide, 7)

paras_wf = [
    ("Design rationale", 14, True, POLIMI_DARK, 6),
    ("• Expanding window: train always starts at t₀ — mimics real deployment", 11, False, DARK_GRAY, 3),
    ("• 4-week embargo: we drop the first 4 observations between train end and val start. This eliminates data leakage from all 4-week lookback features (Δ4w, rolling windows, EWMA warm-up). No temporal overlap is possible.", 11, False, DARK_GRAY, 6),
    ("• Test holdout (2019–2021): SEALED. Never touched during selection or tuning.", 11, True, POLIMI_RED, 8),
    ("Why weighted-by-n_pos aggregation?", 14, True, POLIMI_DARK, 4),
    ("Folds 1–2 have 53/64 positives (GFC, Euro Debt) — statistically reliable. Folds 3–5 have only 2/10/6 positives because 2013–2018 is a structurally low-volatility regime. This is a data property, not a design error.", 11, False, DARK_GRAY, 4),
    ("Weighting by n_pos ensures folds 1–2 dominate the CV metric while folds 3–5 contribute at the margin without introducing noise.", 11, False, DARK_GRAY, 0),
]
add_rich_text(slide, Inches(0.7), Inches(1.3), Inches(6.5), Inches(5.5), paras_wf)

# Fold table as text box
add_shape(slide, Inches(7.8), Inches(1.3), Inches(5), Inches(3.5), LIGHT_GRAY)
paras_ft = [
    ("Fold    Val Period       Y=1%    n_pos   Crisis", 11, True, POLIMI_DARK, 4),
    ("  1      2007–2009      34.9%      53     GFC 2008", 10, False, DARK_GRAY, 2),
    ("  2      2010–2012      41.7%      64     Euro Debt", 10, False, DARK_GRAY, 2),
    ("  3      2013–2014        2.0%        2     Taper 2013", 10, False, MID_GRAY, 2),
    ("  4      2015–2016      10.1%      10     China-Oil", 10, False, MID_GRAY, 2),
    ("  5      2017–2018        6.1%        6     Q4 Selloff", 10, False, MID_GRAY, 6),
    ("Folds 3–5 are thin → n_pos weighting essential", 10, True, POLIMI_RED, 0),
]
add_rich_text(slide, Inches(8), Inches(1.5), Inches(4.6), Inches(3.3), paras_ft)

# Scaler box
add_shape(slide, Inches(7.8), Inches(5.0), Inches(5), Inches(1.8), RGBColor(255, 240, 240))
paras_sc = [
    ("Scaler discipline (critical)", 12, True, POLIMI_RED, 4),
    ("• StandardScaler fit ONLY on each fold's train", 10, False, DARK_GRAY, 2),
    ("• Test holdout: scaler fit on full development set", 10, False, DARK_GRAY, 2),
    ("• AE early stopping: temporal sub-split of train", 10, False, DARK_GRAY, 2),
    ("This is the #1 source of silent leakage in ML finance", 10, True, POLIMI_RED, 0),
]
add_rich_text(slide, Inches(8), Inches(5.1), Inches(4.6), Inches(1.6), paras_sc)


# ════════════════════════════════════════════════
# SLIDE 8: FOUR MODELS
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_bar(slide, "PROCESS & EXECUTION")
add_slide_title(slide, "Four Anomaly Detection Models", "Novelty detection: trained on normal weeks (Y=0) only")
add_logo_corner(slide)
add_footer(slide, 8)

# MVG
add_shape(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.5), RGBColor(240, 248, 255))
paras_mvg = [
    ("1. Multivariate Gaussian (MVG)", 14, True, POLIMI_DARK, 4),
    ("Score = squared Mahalanobis distance: d²(x) = (x−μ)ᵀ Σ⁻¹ (x−μ)", 11, False, DARK_GRAY, 4),
    ("Ledoit-Wolf shrinkage is mandatory: with 56 features and ~360 normals in Fold 1, the sample covariance is near-singular. LW computes the optimal shrinkage Σ* = (1−α)Σ + α·(tr(Σ)/p)·I minimizing expected Frobenius loss — producing a well-conditioned estimator.", 10, False, DARK_GRAY, 0),
]
add_rich_text(slide, Inches(0.7), Inches(1.6), Inches(5.6), Inches(2.3), paras_mvg)

# SVM
add_shape(slide, Inches(0.5), Inches(4.2), Inches(6), Inches(2.5), RGBColor(245, 255, 245))
paras_svm = [
    ("2. One-Class SVM (RBF kernel)", 14, True, POLIMI_DARK, 4),
    ("Maps normals to high-dim feature space, finds smallest hypersphere.", 11, False, DARK_GRAY, 3),
    ("Grid: ν ∈ {0.05, 0.10, 0.15, 0.22} × γ ∈ {scale, auto, 0.01, 0.001}", 11, False, DARK_GRAY, 3),
    ("Best: ν=0.05, γ=0.001 — selected by weighted F1", 11, True, POLIMI_BLUE, 0),
]
add_rich_text(slide, Inches(0.7), Inches(4.3), Inches(5.6), Inches(2.3), paras_svm)

# AE
add_shape(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(2.5), RGBColor(255, 248, 240))
paras_ae = [
    ("3. Autoencoder (symmetric bottleneck)", 14, True, POLIMI_DARK, 4),
    ("56 → 24 → 12 → 6 → 12 → 24 → 56", 12, True, POLIMI_BLUE, 3),
    ("ReLU, dropout 0.15, Adam lr=1e-3, MSE loss", 11, False, DARK_GRAY, 3),
    ("Early stopping on temporal 85/15 sub-split of train normals — NEVER on the CV validation fold (leakage)", 10, False, POLIMI_RED, 0),
]
add_rich_text(slide, Inches(7), Inches(1.6), Inches(5.6), Inches(2.3), paras_ae)

# IF
add_shape(slide, Inches(6.8), Inches(4.2), Inches(6), Inches(2.5), RGBColor(250, 245, 255))
paras_if = [
    ("4. Isolation Forest", 14, True, POLIMI_DARK, 4),
    ("200 trees, random partitioning. Anomalies isolated in fewer splits.", 11, False, DARK_GRAY, 3),
    ("Grid: contamination ∈ {0.05, 0.10, 0.15, 0.22}. Best: 0.10", 11, False, DARK_GRAY, 3),
    ("Threshold fine-tuned via walk-forward weighted by n_pos", 11, False, DARK_GRAY, 0),
]
add_rich_text(slide, Inches(7), Inches(4.3), Inches(5.6), Inches(2.3), paras_if)


# ════════════════════════════════════════════════
# SLIDE 9: MODEL RESULTS
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_bar(slide, "PROCESS & EXECUTION")
add_slide_title(slide, "Model Results — Test Holdout (2019–2021)")
add_logo_corner(slide)
add_footer(slide, 9)

# Results table
add_shape(slide, Inches(0.5), Inches(1.4), Inches(6.2), Inches(2.8), LIGHT_GRAY)
paras_res = [
    ("           F1      AUC-ROC  AUC-PR   Prec    Recall   F2", 11, True, POLIMI_DARK, 4),
    ("MVG     0.516    0.886     0.770    1.000   0.348   0.400", 10, False, DARK_GRAY, 2),
    ("SVM     0.615    0.851     0.732    0.750   0.522   0.556", 10, True, POLIMI_BLUE, 2),
    ("IF       0.583    0.810     0.680    0.560   0.609   0.598", 10, False, DARK_GRAY, 2),
    ("AE      0.516    0.873     0.764    1.000   0.348   0.400", 10, False, DARK_GRAY, 0),
]
add_rich_text(slide, Inches(0.7), Inches(1.6), Inches(5.8), Inches(2.5), paras_res)

paras_obs = [
    ("Key observations:", 13, True, POLIMI_DARK, 4),
    ("• MVG & AE: perfect precision but low recall — conservative thresholds", 11, False, DARK_GRAY, 3),
    ("• SVM: best single model on F1 (balanced precision/recall)", 11, False, DARK_GRAY, 3),
    ("• IF: highest recall (catches 61% of risk-off weeks)", 11, False, DARK_GRAY, 3),
    ("• MVG highest AUC-ROC (0.886): ranking is good, threshold conservative", 11, False, DARK_GRAY, 3),
    ("• F2 (β=2) weights recall double: missing a crisis at 1.5× is costly", 11, True, POLIMI_RED, 0),
]
add_rich_text(slide, Inches(0.5), Inches(4.4), Inches(6.2), Inches(3), paras_obs)

slide.shapes.add_picture(img("fig_mvg_scores.png"), Inches(7), Inches(1.3), Inches(6))


# ════════════════════════════════════════════════
# SLIDE 10: ENSEMBLE
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_bar(slide, "PROCESS & EXECUTION")
add_slide_title(slide, "Ensemble: Combining Four Diverse Detectors")
add_logo_corner(slide)
add_footer(slide, 10)

paras_ens = [
    ("Why ensemble?", 14, True, POLIMI_DARK, 4),
    ("Each model captures different anomaly aspects: distributional shift (MVG), boundary violations (SVM), reconstruction difficulty (AE), structural isolation (IF). If errors are uncorrelated, combining reduces both FP and FN.", 11, False, DARK_GRAY, 8),
    ("Percentile mapping:", 13, True, POLIMI_DARK, 4),
    ("Raw scores are on incomparable scales. We convert each to its empirical percentile against the train-normals distribution: p = F_train(s). All scores become ∈ [0,1].", 11, False, DARK_GRAY, 6),
    ("Three variants:", 13, True, POLIMI_DARK, 4),
    ("1. Hard voting: majority (≥3/4) wins", 11, False, DARK_GRAY, 2),
    ("2. Soft mean: average of percentiles, threshold τ tuned on CV", 11, False, DARK_GRAY, 2),
    ("3. Soft median: robust to one rogue model", 11, False, DARK_GRAY, 6),
    ("Winner: Soft Mean (F1=0.650, beats best single 0.615)", 13, True, POLIMI_GREEN, 0),
]
add_rich_text(slide, Inches(0.7), Inches(1.3), Inches(6), Inches(5.5), paras_ens)

# Table
add_shape(slide, Inches(7), Inches(1.3), Inches(5.8), Inches(2.5), LIGHT_GRAY)
paras_et = [
    ("                  F1       AUC-ROC  AUC-PR", 11, True, POLIMI_DARK, 3),
    ("Hard Vote     0.516    0.859     0.742", 10, False, DARK_GRAY, 2),
    ("Soft Mean     0.650    0.859     0.742", 10, True, POLIMI_GREEN, 2),
    ("Soft Median  0.619    0.869     0.759", 10, False, DARK_GRAY, 0),
]
add_rich_text(slide, Inches(7.2), Inches(1.5), Inches(5.4), Inches(2.3), paras_et)

# Error correlation
slide.shapes.add_picture(img("fig_error_corr.png"), Inches(7.5), Inches(4.0), Inches(4.5))
add_text_box(slide, Inches(7), Inches(6.5), Inches(5.8), Inches(0.5),
             "Mean pairwise error correlation = 0.66 < 0.85 → ensemble adds genuine value",
             font_size=10, bold=True, color=POLIMI_GREEN)


# ════════════════════════════════════════════════
# SLIDE 11: ENSEMBLE TIMELINE
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_bar(slide, "PROCESS & EXECUTION")
add_slide_title(slide, "Ensemble Predictions on Test Holdout")
add_logo_corner(slide)
add_footer(slide, 11)

slide.shapes.add_picture(img("fig_ensemble_timeline.png"), Inches(0.7), Inches(1.4), Inches(11.8), Inches(5))
add_text_box(slide, Inches(0.7), Inches(6.5), Inches(11.8), Inches(0.5),
             "Top: Soft mean percentile score vs. threshold. Red = true risk-off. Bottom: predicted vs. actual signals. The ensemble captures the COVID cluster and correctly filters calm periods.",
             font_size=10, color=MID_GRAY)


# ════════════════════════════════════════════════
# SLIDE 12: BACKTEST
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_bar(slide, "PROCESS & EXECUTION")
add_slide_title(slide, "Backtest Results: Strategy vs. Benchmarks (2019–2021)")
add_logo_corner(slide)
add_footer(slide, 12)

slide.shapes.add_picture(img("fig_equity_curves.png"), Inches(0.3), Inches(1.2), Inches(7.8))

# Metrics table
add_shape(slide, Inches(8.3), Inches(1.2), Inches(4.5), Inches(2.8), LIGHT_GRAY)
paras_bt = [
    ("              Strategy    60/40     B&H", 11, True, POLIMI_DARK, 4),
    ("CAGR         35.2%      18.1%    24.2%", 11, False, DARK_GRAY, 2),
    ("Sharpe       1.93        1.44      1.38", 11, False, DARK_GRAY, 2),
    ("Max DD     −12.1%    −21.5%   −27.7%", 11, True, POLIMI_RED, 2),
    ("Calmar       2.92        0.84      0.87", 11, False, DARK_GRAY, 2),
    ("Turnover    3.9/yr", 11, False, DARK_GRAY, 0),
]
add_rich_text(slide, Inches(8.5), Inches(1.4), Inches(4.1), Inches(2.6), paras_bt)

# Allocation breakdown
add_shape(slide, Inches(8.3), Inches(4.2), Inches(4.5), Inches(1.8), RGBColor(240, 255, 240))
paras_al = [
    ("Allocation breakdown:", 12, True, POLIMI_DARK, 4),
    ("Levered Equity: 103 weeks (85.8%)", 11, False, POLIMI_GREEN, 2),
    ("Gold: 8 weeks (6.7%)", 11, False, RGBColor(200, 160, 0), 2),
    ("Cash USD: 8 weeks (6.7%)", 11, False, POLIMI_BLUE, 2),
    ("MBS: 1 week (0.8%)", 11, False, MID_GRAY, 0),
]
add_rich_text(slide, Inches(8.5), Inches(4.3), Inches(4.1), Inches(1.6), paras_al)

paras_ins = [
    ("Leverage generates alpha in calm periods (103w at 1.5×),", 10, True, POLIMI_DARK, 2),
    ("while EWS limits drawdowns: −12% vs −28% buy&hold.", 10, True, POLIMI_DARK, 2),
    ("Leverage + early warning = superior Sharpe ratio.", 10, True, POLIMI_GREEN, 0),
]
add_rich_text(slide, Inches(8.3), Inches(6.2), Inches(4.5), Inches(1), paras_ins)


# ════════════════════════════════════════════════
# SLIDE 13: COVID STRESS TEST
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_bar(slide, "PROCESS & EXECUTION")
add_slide_title(slide, "COVID Stress Test: Trigger Analysis & Allocation Routing")
add_logo_corner(slide)
add_footer(slide, 13)

slide.shapes.add_picture(img("fig_stress_timeline.png"), Inches(0.3), Inches(1.2), Inches(7.5), Inches(3))
slide.shapes.add_picture(img("fig_covid_triggers.png"), Inches(0.3), Inches(4.3), Inches(7.5), Inches(3))

paras_covid = [
    ("COVID Crash (Feb–Apr 2020)", 13, True, POLIMI_RED, 4),
    ("• Ensemble signals risk-off as VIX spikes", 10, False, DARK_GRAY, 2),
    ("• USD triggers fire first: funding stress (LIBOR-3M widens), DXY rallies, VRP explodes", 10, False, DARK_GRAY, 2),
    ("• Routing → Cash USD and Gold", 10, False, DARK_GRAY, 2),
    ("• MBS correctly blocked (VIX > 30)", 10, False, DARK_GRAY, 8),
    ("COVID Recovery (Apr–Dec 2020)", 13, True, POLIMI_GREEN, 4),
    ("• Risk-off signals fade as markets stabilize", 10, False, DARK_GRAY, 2),
    ("• Returns to 1.5× equity, captures V-shaped rebound", 10, False, DARK_GRAY, 2),
    ("• Gold score stays elevated (fiscal stimulus → inflation expectations)", 10, False, DARK_GRAY, 8),
    ("Reflation 2021", 13, True, POLIMI_BLUE, 4),
    ("• Full equity exposure maintained", 10, False, DARK_GRAY, 2),
    ("• No false alarms during re-opening rally", 10, False, DARK_GRAY, 0),
]
add_rich_text(slide, Inches(8), Inches(1.2), Inches(5), Inches(6), paras_covid)


# ════════════════════════════════════════════════
# SLIDE 14: SUB-SCORES
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_bar(slide, "PROCESS & EXECUTION")
add_slide_title(slide, "Domain Sub-Score Dynamics on Test Holdout")
add_logo_corner(slide)
add_footer(slide, 14)

slide.shapes.add_picture(img("fig_subscores.png"), Inches(0.5), Inches(1.2), Inches(8))

paras_sub = [
    ("Interpretation:", 14, True, POLIMI_DARK, 6),
    ("USD sub-score (top): spikes during COVID crash — funding stress, DXY rally, VRP explosion", 11, False, POLIMI_BLUE, 4),
    ("Gold sub-score (mid): elevated during COVID and recovery — stimulus drives low real yields, DXY weakens", 11, False, RGBColor(200, 160, 0), 4),
    ("MBS (bottom): activates only once — strict conditions correctly filter out acute stress", 11, False, MID_GRAY, 8),
    ("Red shading = true risk-off weeks. Sub-scores provide interpretable routing: you can trace which macro regime drove each allocation decision.", 11, False, DARK_GRAY, 0),
]
add_rich_text(slide, Inches(8.7), Inches(1.3), Inches(4.2), Inches(5.5), paras_sub)


# ════════════════════════════════════════════════
# SLIDE 15: LIMITATIONS & CONCLUSIONS
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_bar(slide, "PROCESS & EXECUTION")
add_slide_title(slide, "Limitations, Conclusions & Future Work")
add_logo_corner(slide)
add_footer(slide, 15)

# Limitations
add_shape(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(3.2), RGBColor(255, 245, 245))
paras_lim = [
    ("Limitations (candidly)", 14, True, POLIMI_RED, 6),
    ("• Short test: 120 weeks, 9 switches — results depend on few trades", 11, False, DARK_GRAY, 3),
    ("• Thin CV folds 3–5: 2–10 positives, n_pos weighting mitigates but doesn't eliminate noise", 11, False, DARK_GRAY, 3),
    ("• Proxy constraints: MXUS ≠ MSCI World; log-ratios ≠ OAS", 11, False, DARK_GRAY, 3),
    ("• Weekly frequency misses intra-week volatility events", 11, False, DARK_GRAY, 3),
    ("• No 2022 data: inflation/rate shock regime not tested", 11, False, DARK_GRAY, 3),
    ("• At 1.5×, recall gap (35%) is a genuine concern", 11, False, DARK_GRAY, 0),
]
add_rich_text(slide, Inches(0.7), Inches(1.5), Inches(5.4), Inches(3), paras_lim)

# Conclusions
add_shape(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(3.2), RGBColor(240, 255, 240))
paras_conc = [
    ("Conclusions", 14, True, POLIMI_GREEN, 6),
    ("• Ensemble (F1=0.65) outperforms all single models", 11, False, DARK_GRAY, 3),
    ("• Domain-specific routing adds genuine intelligence vs. binary risk-on/off", 11, False, DARK_GRAY, 3),
    ("• Walk-forward + embargo = methodologically sound OOS evaluation", 11, False, DARK_GRAY, 3),
    ("• Strategy: Sharpe 1.93, Max DD −12% vs −28% buy&hold", 11, True, POLIMI_DARK, 3),
    ("• Leverage + early warning = superior risk-adjusted returns", 11, True, POLIMI_GREEN, 0),
]
add_rich_text(slide, Inches(7), Inches(1.5), Inches(5.6), Inches(3), paras_conc)

# Future work
add_shape(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2), LIGHT_GRAY)
paras_fw = [
    ("Future Work", 14, True, POLIMI_DARK, 4),
    ("• Daily frequency for finer granularity    • SHAP / Mahalanobis decomposition for signal attribution    • Include 2022–2025 data (inflation regime validation)", 11, False, DARK_GRAY, 3),
    ("• Regime-switching models (HMM) as benchmark    • Bootstrap CV for confidence intervals    • Dynamic threshold adaptation", 11, False, DARK_GRAY, 0),
]
add_rich_text(slide, Inches(0.7), Inches(4.9), Inches(11.9), Inches(2), paras_fw)


# ── SAVE ──
output_path = "presentation/EWS_PoliMi_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Slides: {len(prs.slides)}")
